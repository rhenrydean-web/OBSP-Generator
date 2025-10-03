# app.py
import io, os, json
import streamlit as st
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# ---------- PAGE CONFIG (must be first Streamlit call) ----------
st.set_page_config(page_title="OBSP Generator (Robust)", page_icon="📘", layout="centered")

PRIMARY = RGBColor(30, 64, 175)   # deep blue
BODY_SIZE = 18
TITLE_SIZE = 36

# ---------- UI HELPERS ----------
def add_title(slide, text):
    title = slide.shapes.title
    title.text = text
    p = title.text_frame.paragraphs[0]
    p.font.size = Pt(TITLE_SIZE); p.font.bold = True; p.font.color.rgb = PRIMARY

def add_body(slide, text):
    ph = slide.placeholders[1]
    ph.text = text
    for p in ph.text_frame.paragraphs:
        p.font.size = Pt(BODY_SIZE)

def add_bullets(slide, lines):
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate([l for l in lines if l.strip()]):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.text = line.strip(); p.level = 0; p.font.size = Pt(BODY_SIZE)

def add_two_col_table(slide, left_title, left_items, right_title, right_items):
    from pptx.util import Pt
    rows = max(len(left_items), len(right_items)) + 1
    table = slide.shapes.add_table(rows, 2, Pt(20), Pt(120), Pt(900), Pt(360)).table
    table.cell(0,0).text, table.cell(0,1).text = left_title, right_title
    for j in (0,1):
        p = table.cell(0,j).text_frame.paragraphs[0]
        p.font.bold = True; p.font.size = Pt(BODY_SIZE)
    for i in range(1, rows):
        table.cell(i,0).text = left_items[i-1] if i-1 < len(left_items) else ""
        table.cell(i,1).text = right_items[i-1] if i-1 < len(right_items) else ""

def add_kpi_table(slide, kpis):
    from pptx.util import Pt
    rows, cols = len(kpis)+1, 5
    table = slide.shapes.add_table(rows, cols, Pt(20), Pt(120), Pt(900), Pt(360)).table
    headers = ["Metric","Baseline","Target","Cadence","Owner"]
    for j,h in enumerate(headers):
        c = table.cell(0,j); c.text = h
        p = c.text_frame.paragraphs[0]; p.font.bold = True; p.font.size = Pt(BODY_SIZE)
    for i,k in enumerate(kpis, start=1):
        table.cell(i,0).text = k.get("metric","")
        table.cell(i,1).text = k.get("baseline","")
        table.cell(i,2).text = k.get("target","")
        table.cell(i,3).text = k.get("cadence","")
        table.cell(i,4).text = k.get("owner","")

def extract_json(raw: str):
    start = raw.find("{"); end = raw.rfind("}")
    if start == -1 or end == -1:
        raise ValueError("No JSON object found in LLM response.")
    return json.loads(raw[start:end+1])

def call_llm(system, user, model):
    if "OPENAI_API_KEY" not in st.secrets or not st.secrets["OPENAI_API_KEY"]:
        raise RuntimeError("Missing OPENAI_API_KEY secret. Add it in Streamlit Cloud → App → Settings → Secrets.")
    from openai import OpenAI
    client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])
    resp = client.chat.completions.create(
        model=model,
        messages=[{"role":"system","content":system},{"role":"user","content":user}],
        temperature=0.2,
    )
    return resp.choices[0].message.content

# ---------- Context extraction (lightweight & fast) ----------
from pypdf import PdfReader
from docx import Document as DocxDocument
from rapidfuzz import fuzz

def extract_text_from_file(f):
    name = (f.name or "").lower()
    if name.endswith(".pdf"):
        reader = PdfReader(f)
        pages = min(len(reader.pages), 25)
        return "\n\n".join((reader.pages[i].extract_text() or "") for i in range(pages))
    elif name.endswith(".docx"):
        d = DocxDocument(f)
        return "\n".join(p.text for p in d.paragraphs)
    else:  # .txt
        return f.read().decode(errors="ignore")

def chunk_text(txt, chunk_size=800, overlap=120):
    txt = " ".join(txt.split())
    chunks, i = [], 0
    while i < len(txt):
        chunks.append(txt[i:i+chunk_size])
        i += chunk_size - overlap
    return chunks

def pick_top_snippets(texts, query, k=8):
    chunks = []
    for t in texts:
        chunks.extend(chunk_text(t))
    if not chunks:
        return []
    scored = [(c, fuzz.partial_ratio(query, c)) for c in chunks]
    scored.sort(key=lambda x: x[1], reverse=True)
    return [c for c, s in scored[:k] if s > 40]

# ---------- APP UI ----------
st.markdown("## 📘 Outcome-Based Customer Success Plan — Robust Generator")
st.caption("Fill the form, (optionally) upload context + a PPT template, and download a polished deck.")

with st.form("obsp"):
    col1, col2 = st.columns(2)
    customer = col1.text_input("Customer", "Acme Corp")
    industry = col2.text_input("Industry", "SaaS - HR Tech")
    arr = col1.text_input("ARR Segment", "$1–5M ARR")
    horizon = col2.text_input("Time Horizon", "12 months")

    objectives = st.text_area("Customer Outcomes (comma-separated)", "Reduce churn by 10%, Improve onboarding time by 20%")
    baselines  = st.text_area("Baselines (e.g., NPS=45, Adoption=60%)", "NPS=45, Adoption=60%")
    constraints = st.text_area("Constraints", "Limited customer resources; integration backlog")
    stakeholders = st.text_area("Stakeholders", "Customer: Exec Sponsor, Admin, Champions; Vendor: CSM, SE, Support")
    risks = st.text_area("Known Risks", "Low admin bandwidth; End-user training gaps")

    uploads = st.file_uploader(
        "Upload context files (QBRs, notes, SOWs) — PDF, DOCX, or TXT (optional)",
        type=["pdf","docx","txt"], accept_multiple_files=True
    )
    ppt_template = st.file_uploader(
        "Optional: Upload a PowerPoint template (.pptx) to apply your brand/theme",
        type=["pptx"], accept_multiple_files=False
    )

    model = st.selectbox("Model", ["gpt-4o-mini","gpt-4o","gpt-4-turbo","gpt-3.5-turbo"], index=0)
    submitted = st.form_submit_button("Generate Plan")

if submitted:
    # Build a simple relevance query from the form
    query = f"""
    {customer} {industry} {arr} {horizon}
    Objectives: {objectives}
    Baselines: {baselines}
    Constraints: {constraints}
    Stakeholders: {stakeholders}
    Risks: {risks}
    """
    all_texts = []
    for f in (uploads or []):
        try:
            all_texts.append(extract_text_from_file(f))
        except Exception:
            pass
    snippets = pick_top_snippets(all_texts, query, k=8)
    context_blob = "\n\n---\n\n".join(snippets)[:8000]

    # --------- ROBUST PLAN SCHEMA (aligned to your detailed brief) ----------
    system = "You output STRICT JSON for a comprehensive customer success plan. No markdown, no extra prose."
    user = f"""
Return ONLY JSON with this schema (no markdown):

{{
  "cover": {{
    "title": "Outcome-Based Customer Success Plan — {customer}",
    "account_meta": "Industry: {industry} | ARR: {arr} | Horizon: {horizon}"
  }},
  "purpose": "1-2 short paragraphs: explain the purpose of a customer success plan, focusing on alignment to customer goals.",
  "cs_vs_service": {{
    "summary": "2 short paragraphs contrasting proactive customer success vs reactive customer service.",
    "success_characteristics": ["Proactive","Relationship-driven","Outcome-oriented","Value-focused"],
    "service_characteristics": ["Reactive","Issue resolution","Transactional","Support-focused"]
  }},
  "importance": ["Reduce churn","Increase retention revenue","Enable cross-sell/upsell","Provide customer insights","Lower support costs","Create competitive advantage"],
  "strategy_elements": {{
    "profiles": "Brief summary of customer segments/profiles relevant to this account.",
    "goals_kpis": ["List key KPIs to track (e.g., NPS, adoption, renewal rate, MRR)"],
    "milestones_touchpoints": ["Key milestones and touchpoints in the journey"],
    "tasks_by_stage": ["Key tasks per stage"],
    "prioritization": "How opportunities are prioritized",
    "replicable_processes": ["High-impact repeatable processes"],
    "handoff_guidelines": "Handoff guidelines across teams"
  }},
  "nine_step_guide": {{
    "roles_responsibilities": ["CCO","VP CS","Director CS","CSM","CS Associate"],
    "journey_map_stages": ["Awareness","Consideration","Conversion","Adoption","Advocacy"],
    "success_by_stage": ["Conversion success…","Retention success…","Advocacy success…"],
    "data_for_opportunities": ["Interaction history","CES","Agent performance","Feedback","Support requests"],
    "processes_by_moment": {{
      "awareness": ["What marketing does…"],
      "consideration": ["What sales does…"],
      "conversion": ["Onboarding steps…"],
      "adoption": ["Training, configuration, support, expansion"],
      "advocacy": ["Referrals, reviews, check-ins"]
    }},
    "metrics_and_benchmarks": ["Retention","CLV","NPS","Churn","CES","MRR"],
    "cross_functional_collab": ["Sales↔CS","Marketing↔CS","Support↔CS"],
    "tooling": ["CS platform","Feedback tools","Automation"],
    "feedback_loop": ["Surveys","VOC","Usage telemetry","QBR inputs"]
  }},
  "three_examples": {{
    "onboarding_plan": {{
      "summary": "Short overview for adoption phase",
      "sections": ["Summary","Journey map","Strategy","Product adoption","Resource allocation"]
    }},
    "growth_expansion_plan": {{
      "summary": "Deepen value & expansion",
      "sections": ["Summary","Onboarding foundations","Sales/marketing engagement","Success engagement","Reporting","Revenue generators"]
    }},
    "renewal_plan": {{
      "summary": "Assess past; plan next cycle",
      "sections": ["Overview","Customer overview","Performance","Journey tasks","Opportunities","Roles & responsibilities"]
    }}
  }},
  "best_practices": [
    "Set SMART goals",
    "Segment strategies by journey stage and need",
    "Clarify ownership and outcomes",
    "Automate repeatable tasks",
    "Build trust before selling",
    "Share data and collaborate",
    "Design a clear customer experience"
  ],
  "obsp_core": {{
    "objectives": [{", ".join([f'"{o.strip()}"' for o in objectives.split(",")])}],
    "kpis": [
      {{"metric":"NPS","baseline":"{baselines}","target":"Improve vs baseline","cadence":"Monthly","owner":"CSM"}}
    ],
    "milestones": [
      {{"phase":"Onboarding (0–30 days)","deliverables":["Training complete","Integrations configured"]}},
      {{"phase":"Adoption (30–90 days)","deliverables":[">60% WAU","Feature A rolled out"]}},
      {{"phase":"Optimization (90–180 days)","deliverables":["Advanced features live","Process improvements"]}},
      {{"phase":"Renewal Prep (180–365 days)","deliverables":["ROI case","Exec QBR"]}}
    ],
    "roles_vendor": ["CSM — governance & outcomes","SE — integrations","Support — issue resolution"],
    "roles_customer": ["Exec Sponsor — vision & unblock","Admin — configuration","Champions — adoption"]
  }},
  "governance": "Narrative: weekly status, monthly steering, quarterly business reviews, clear escalation path."
}}

Context about the specific account:
Customer: {customer}
Industry: {industry}
ARR Segment: {arr}
Outcomes: {objectives}
Baselines: {baselines}
Horizon: {horizon}
Constraints: {constraints}
Stakeholders: {stakeholders}
Known Risks: {risks}

Relevant Context (use when helpful; paraphrase succinctly, cite in-line): 
{context_blob}
"""
    try:
        raw = call_llm(system, user, model)
        data = extract_json(raw)
    except Exception as e:
        st.error(f"Could not create plan. {e}")
        st.stop()

    # ---------- Build PPT (use uploaded template if provided) ----------
    prs = Presentation(ppt_template) if ppt_template is not None else Presentation()

    # Title
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s, data["cover"]["title"])
    s.placeholders[1].text = data["cover"]["account_meta"]
    for p in s.placeholders[1].text_frame.paragraphs: p.font.size = Pt(18)

    # Purpose
    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "1. Purpose of a Customer Success Plan")
    add_body(s, data.get("purpose",""))

    # CS vs Service
    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "2. Customer Success vs Customer Service")
    body = data.get("cs_vs_service", {})
    text = body.get("summary","")
    add_body(s, text)

    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "2a. Key Characteristics")
    left = body.get("success_characteristics", [])
    right = body.get("service_characteristics", [])
    add_two_col_table(s, "Customer Success (Proactive)", left, "Customer Service (Reactive)", right)

    # Importance
    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "3. Why It Matters")
    add_bullets(s, data.get("importance", []))

    # Strategy elements
    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "4. Plan Elements")
    se = data.get("strategy_elements", {})
    se_lines = [
        f"Profiles: {se.get('profiles','')}",
        f"Goals & KPIs: {', '.join(se.get('goals_kpis', []))}",
        f"Milestones/Touchpoints: {', '.join(se.get('milestones_touchpoints', []))}",
        f"Tasks by Stage: {', '.join(se.get('tasks_by_stage', []))}",
        f"Prioritization: {se.get('prioritization','')}",
        f"Replicable processes: {', '.join(se.get('replicable_processes', []))}",
        f"Handoff guidelines: {se.get('handoff_guidelines','')}",
    ]
    add_bullets(s, se_lines)

    # Nine-step guide (condensed across a few slides)
    ns = data.get("nine_step_guide", {})

    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "5. Roles & Journey")
    add_two_col_table(s, "Roles/Responsibilities", ns.get("roles_responsibilities", []),
                         "Journey Stages", ns.get("journey_map_stages", []))

    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "6. Success by Stage")
    add_bullets(s, ns.get("success_by_stage", []))

    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "7. Data → High-Impact Opportunities")
    add_bullets(s, ns.get("data_for_opportunities", []))

    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "8. Processes by Moment in Journey")
    pbm = ns.get("processes_by_moment", {})
    lines = []
    for k in ["awareness","consideration","conversion","adoption","advocacy"]:
        if k in pbm:
            lines.append(f"{k.title()}: " + "; ".join(pbm[k]))
    add_bullets(s, lines)

    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "9. Metrics & Benchmarks")
    add_bullets(s, ns.get("metrics_and_benchmarks", []))

    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "10. Cross-Functional Collaboration & Tooling")
    add_two_col_table(s, "Collaboration", ns.get("cross_functional_collab", []),
                         "Tooling", ns.get("tooling", []))

    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "11. Feedback Loop")
    add_bullets(s, ns.get("feedback_loop", []))

    # Three plan examples
    ex = data.get("three_examples", {})
    for title_key, human_title in [("onboarding_plan","Example Plan: Onboarding"),
                                   ("growth_expansion_plan","Example Plan: Growth & Expansion"),
                                   ("renewal_plan","Example Plan: Renewal")]:
        if title_key in ex:
            s = prs.slides.add_slide(prs.slide_layouts[1])
            add_title(s, human_title)
            x = ex[title_key]
            add_body(s, x.get("summary",""))
            s2 = prs.slides.add_slide(prs.slide_layouts[1])
            add_title(s2, f"{human_title} — Sections")
            add_bullets(s2, x.get("sections", []))

    # Best practices
    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "Best Practices")
    add_bullets(s, data.get("best_practices", []))

    # OBSP core (objectives, KPIs, milestones, roles) + governance
    core = data.get("obsp_core", {})
    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "Account Objectives")
    add_bullets(s, core.get("objectives", []))

    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "Success Metrics & KPIs")
    kpis = core.get("kpis", [])
    from pptx.util import Pt
    if kpis:
        from pptx.util import Pt
        rows, cols = len(kpis)+1, 5
        table = s.shapes.add_table(rows, cols, Pt(20), Pt(120), Pt(900), Pt(360)).table
        headers = ["Metric","Baseline","Target","Cadence","Owner"]
        for j,h in enumerate(headers):
            c = table.cell(0,j); c.text = h
            p = c.text_frame.paragraphs[0]; p.font.bold = True; p.font.size = Pt(BODY_SIZE)
        for i,k in enumerate(kpis, start=1):
            table.cell(i,0).text = k.get("metric","")
            table.cell(i,1).text = k.get("baseline","")
            table.cell(i,2).text = k.get("target","")
            table.cell(i,3).text = k.get("cadence","")
            table.cell(i,4).text = k.get("owner","")
    else:
        add_body(s, "No KPIs provided.")

    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "Milestones & Timeline")
    lines = [f"{m.get('phase','')}: " + "; ".join(m.get('deliverables',[])) for m in core.get("milestones",[])]
    add_bullets(s, lines)

    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "Roles & Responsibilities")
    add_two_col_table(s, "Vendor Team", core.get("roles_vendor", []), "Customer Team", core.get("roles_customer", []))

    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(s, "Engagement & Governance")
    add_body(s, data.get("governance", ""))

    # Export
    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    st.success("Plan generated!")
    st.download_button(
        "⬇️ Download PowerPoint",
        buf,
        file_name=f"OBSP_{customer.replace(' ','_')}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
